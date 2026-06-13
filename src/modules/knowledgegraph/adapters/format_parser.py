"""
Format Parser — extracts text content from diverse file formats for knowledge extraction.

Handles: plain text, markdown, CSV, JSON, log files, Word (.docx), PDF, Excel (.xlsx), PowerPoint (.pptx)
All content is normalized to markdown-like text for downstream KnowledgeExtractor.

No audio/video support (out of scope).

Standards:
    - CODDY-Architecture-v1.0    → Adapter Pattern (wraps python-docx, pypdf, openpyxl, python-pptx)
    - CODDY-ProjectStructure-v1.0 → adapters/ requirements
"""

from __future__ import annotations

__all__ = ["FormatParser"]

import csv
import io
import json
import logging
import re
from pathlib import Path
from typing import Dict, List, Optional, Tuple

logger = logging.getLogger("CodeCortex.KnowledgeGraph.FormatParser")


class FormatParser:
    """Extract normalized text from any supported file format."""

    # Extensions we can handle (lowercase, with dot)
    SUPPORTED = {
        ".md", ".rst", ".txt", ".adoc",       # text/markup
        ".csv", ".json", ".log",             # structured text
        ".docx",                              # Word
        ".pdf",                              # PDF
        ".xlsx", ".xls",                     # Excel
        ".pptx", ".ppt",                     # PowerPoint
    }

    @classmethod
    def can_parse(cls, file_path: str) -> bool:
        return Path(file_path).suffix.lower() in cls.SUPPORTED

    @classmethod
    def extract(cls, file_path: str) -> Tuple[str, Optional[str]]:
        """Extract normalized text from file.

        Returns:
            (content, error) where content is markdown-like text,
            error is None on success or error message on failure.
        """
        path = Path(file_path)
        ext = path.suffix.lower()

        if ext in (".md", ".rst", ".txt", ".adoc"):
            return cls._parse_text(path)
        if ext == ".csv":
            return cls._parse_csv(path)
        if ext == ".json":
            return cls._parse_json(path)
        if ext == ".log":
            return cls._parse_log(path)
        if ext == ".docx":
            return cls._parse_docx(path)
        if ext == ".pdf":
            return cls._parse_pdf(path)
        if ext in (".xlsx", ".xls"):
            return cls._parse_excel(path)
        if ext in (".pptx", ".ppt"):
            return cls._parse_pptx(path)

        return "", f"Unsupported format: {ext}"

    # ── Plain text / markup ─────────────────────────────────

    @staticmethod
    def _parse_text(path: Path) -> Tuple[str, Optional[str]]:
        try:
            content = path.read_text(encoding="utf-8", errors="replace")
            return content, None
        except Exception as e:
            return "", str(e)

    # ── CSV ─────────────────────────────────────────────────

    @staticmethod
    def _parse_csv(path: Path) -> Tuple[str, Optional[str]]:
        try:
            lines: List[str] = ["# CSV Data Extract\n"]
            with path.open("r", encoding="utf-8", errors="replace", newline="") as f:
                reader = csv.reader(f)
                headers = next(reader, None)
                if headers:
                    lines.append("## Columns\n")
                    for h in headers:
                        lines.append(f"- {h}")
                    lines.append("")
                lines.append("## Rows\n")
                for i, row in enumerate(reader, 1):
                    if i > 500:
                        lines.append("\n... (truncated at 500 rows)")
                        break
                    row_text = " | ".join(f"{headers[j] if headers else j}: {v}" for j, v in enumerate(row) if v)
                    if row_text:
                        lines.append(f"- Row {i}: {row_text}")
            return "\n".join(lines), None
        except Exception as e:
            return "", f"CSV parse error: {e}"

    # ── JSON ────────────────────────────────────────────────

    @staticmethod
    def _parse_json(path: Path) -> Tuple[str, Optional[str]]:
        try:
            content = path.read_text(encoding="utf-8", errors="replace")
            data = json.loads(content)
            lines = ["# JSON Data Extract\n"]
            lines.extend(FormatParser._flatten_json(data, max_depth=4))
            return "\n".join(lines), None
        except json.JSONDecodeError as e:
            # Still return raw text if JSON is invalid
            return content, f"JSON parse warning: {e}"
        except Exception as e:
            return "", f"JSON read error: {e}"

    @staticmethod
    def _flatten_json(data, prefix: str = "", depth: int = 0, max_depth: int = 4) -> List[str]:
        """Flatten JSON to markdown-like key-value lines."""
        lines: List[str] = []
        if depth > max_depth:
            return lines
        if isinstance(data, dict):
            for k, v in data.items():
                key_path = f"{prefix}.{k}" if prefix else k
                if isinstance(v, (dict, list)) and depth < max_depth:
                    lines.append(f"## {key_path}")
                    lines.extend(FormatParser._flatten_json(v, key_path, depth + 1, max_depth))
                else:
                    val_str = str(v)[:200]
                    lines.append(f"- **{key_path}**: {val_str}")
        elif isinstance(data, list):
            for i, item in enumerate(data[:50]):  # cap list items
                item_prefix = f"{prefix}[{i}]"
                if isinstance(item, (dict, list)) and depth < max_depth:
                    lines.append(f"## {item_prefix}")
                    lines.extend(FormatParser._flatten_json(item, item_prefix, depth + 1, max_depth))
                else:
                    val_str = str(item)[:200]
                    lines.append(f"- **{item_prefix}**: {val_str}")
        else:
            lines.append(f"- **{prefix}**: {str(data)[:200]}")
        return lines

    # ── Log files ───────────────────────────────────────────

    @staticmethod
    def _parse_log(path: Path) -> Tuple[str, Optional[str]]:
        try:
            content = path.read_text(encoding="utf-8", errors="replace")
            lines = content.splitlines()
            # Extract ERROR, WARN, CRITICAL lines as knowledge
            knowledge_lines = ["# Log Analysis\n"]
            error_pattern = re.compile(r"\b(ERROR|CRITICAL|FATAL|EXCEPTION|WARN|WARNING|DEBUG|INFO)\b", re.I)
            seen: set = set()
            for line in lines[:2000]:  # cap lines
                m = error_pattern.search(line)
                if m:
                    key = line.strip()[:120]
                    if key not in seen:
                        seen.add(key)
                        knowledge_lines.append(f"- {key}")
            if len(knowledge_lines) == 1:
                knowledge_lines.append("\n## Sample Lines\n")
                for line in lines[:20]:
                    knowledge_lines.append(f"- {line[:200]}")
            return "\n".join(knowledge_lines), None
        except Exception as e:
            return "", f"Log parse error: {e}"

    # ── Word (.docx) ────────────────────────────────────────

    @staticmethod
    def _parse_docx(path: Path) -> Tuple[str, Optional[str]]:
        try:
            from docx import Document
            doc = Document(str(path))
            lines: List[str] = ["# Word Document Extract\n"]
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    lines.append(text)
            # Tables
            for table in doc.tables[:20]:
                lines.append("\n## Table\n")
                for row in table.rows[:50]:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text:
                        lines.append(f"- {row_text}")
            return "\n".join(lines), None
        except ImportError:
            return "", "python-docx not installed"
        except Exception as e:
            return "", f"DOCX parse error: {e}"

    # ── PDF ─────────────────────────────────────────────────

    @staticmethod
    def _parse_pdf(path: Path) -> Tuple[str, Optional[str]]:
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            lines: List[str] = ["# PDF Extract\n"]
            for i, page in enumerate(reader.pages[:50]):  # cap pages
                text = page.extract_text()
                if text:
                    lines.append(f"\n## Page {i + 1}\n")
                    lines.append(text)
            return "\n".join(lines), None
        except ImportError:
            return "", "pypdf not installed"
        except Exception as e:
            return "", f"PDF parse error: {e}"

    # ── Excel (.xlsx / .xls) ────────────────────────────────

    @staticmethod
    def _parse_excel(path: Path) -> Tuple[str, Optional[str]]:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(path), data_only=True, read_only=True)
            lines: List[str] = [f"# Excel Extract: {path.name}\n"]
            for sheet_name in wb.sheetnames[:10]:
                ws = wb[sheet_name]
                lines.append(f"\n## Sheet: {sheet_name}\n")
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    row_count += 1
                    if row_count > 200:
                        lines.append("... (truncated)")
                        break
                    row_text = " | ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        lines.append(f"- {row_text}")
            return "\n".join(lines), None
        except ImportError:
            return "", "openpyxl not installed"
        except Exception as e:
            return "", f"Excel parse error: {e}"

    # ── PowerPoint (.pptx / .ppt) ───────────────────────────

    @staticmethod
    def _parse_pptx(path: Path) -> Tuple[str, Optional[str]]:
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            lines: List[str] = [f"# PowerPoint Extract: {path.name}\n"]
            for i, slide in enumerate(prs.slides[:50], 1):
                lines.append(f"\n## Slide {i}\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
                    # Tables in slides
                    if shape.has_table:
                        table = shape.table
                        lines.append("\n### Table\n")
                        for row in table.rows[:20]:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            if row_text:
                                lines.append(f"- {row_text}")
            return "\n".join(lines), None
        except ImportError:
            return "", "python-pptx not installed"
        except Exception as e:
            return "", f"PPTX parse error: {e}"
