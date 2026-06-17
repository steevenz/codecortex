"""
Unified Search Orchestrator — 17 providers with 9Router-compatible API.

Providers:
  codecortex-codebase      — FTS5 + semantic + graph symbol search
  codecortex-repowt        — Git working-tree, status, commits, diffs
  codecortex-filesystem    — File glob + content regex + ReDoS-safe
  codecortex-graph         — Graph relationship/trace/hierarchy search
  codecortex-idegraph      — Cross-IDE memory/conversation search
  codecortex-knowledge     — Knowledge graph FTS + structured queries
  codecortex-crossproject  — Cross-project reference tracking
  codecortex-codeindex     — Fast code index metadata lookup
  codecortex-agentart      — .agents folder artifact search
  codecortex-codelogs      — Log file search with level/time filters
  codecortex-todo          — Comment tag scanner (TODO, FIXME, HACK, XXX, STUB, etc.)
  codecortex-stub          — Empty function/class stub detection
  codecortex-security      — Security audit: secrets, vulns, PII, misconfig
  codecortex-empty         — Empty file and folder detection
  codecortex-svn           — SVN working tree status, log, info
  codecortex-blame         — Git blame, hotspots, commit history
  codecortex-combo         — All providers orchestrated (default)

:project: CodeCortex
:package: Services.UnifiedSearch
:author: Steeven Andrian
:copyright: (c) 2026 CODDY Codework
:standard: CODDY-API-v1.0
"""
from __future__ import annotations
import asyncio
import os
import re
import time
import glob as globmod
import fnmatch
import logging
import subprocess
import threading
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger("CodeCortex.UnifiedSearch")

# ────────────────────────────────────────────────────────────
# ReDoS Protection
# ────────────────────────────────────────────────────────────
REGEX_TIMEOUT_SECONDS = 5
REGEX_MAX_INPUT_BYTES = 500 * 1024
REGEX_MAX_BACKTRACK = 100_000

_regex_timeout_lock = threading.Lock()

def _safe_regex_compile(pattern: str, flags: int = re.IGNORECASE) -> re.Pattern:
    """Compile regex with ReDoS timeout & complexity guard."""
    if len(pattern) > 2000:
        raise ValueError(f"Regex pattern too long ({len(pattern)} chars, max 2000)")
    dangerous = [
        r'(.+)+', r'([a-zA-Z]+)*', r'(a+)+', r'(.+){', r'.*.*.*',
    ]
    for d in dangerous:
        if d in pattern:
            raise ValueError(f"Potentially catastrophic regex pattern: {d}")
    try:
        return re.compile(pattern, flags)
    except re.error as e:
        raise ValueError(f"Invalid regex: {e}")

def _safe_regex_search(pattern: re.Pattern, text: str, max_chars: int = REGEX_MAX_INPUT_BYTES) -> List:
    """Execute regex with input size limit (ReDoS prevention)."""
    if len(text) > max_chars:
        text = text[:max_chars]
    start = time.monotonic()
    try:
        results = list(pattern.finditer(text))
    except Exception as e:
        logger.warning("Regex search error: %s", e)
        return []
    elapsed = time.monotonic() - start
    if elapsed > REGEX_TIMEOUT_SECONDS:
        logger.warning("Regex search exceeded timeout: %.2fs", elapsed)
    return results

# ────────────────────────────────────────────────────────────
# Path Validation
# ────────────────────────────────────────────────────────────
def _validate_path(search_path: str, allowed_roots: Optional[List[str]] = None) -> Path:
    """Validate and resolve path, preventing traversal outside allowed roots."""
    try:
        resolved = Path(search_path).resolve()
    except Exception:
        raise ValueError(f"Invalid path: {search_path}")
    if not allowed_roots:
        allowed_roots = [str(Path.home()), str(Path.cwd())]
    root_resolved = [str(Path(r).resolve()) for r in allowed_roots if os.path.exists(r)]
    ok = any(str(resolved).startswith(r) for r in root_resolved)
    if not ok:
        raise ValueError(f"Path traversal denied: {search_path}")
    return resolved


def _extract_document_text(path: str) -> Optional[str]:
    """Extract text from PDF, DOCX, XLSX binary documents."""
    ext = Path(path).suffix.lower()
    try:
        if ext == ".pdf":
            try:
                from pypdf import PdfReader
                reader = PdfReader(path)
                pages = [p.extract_text() for p in reader.pages if p.extract_text() and p.extract_text().strip()]
                return "\n".join(pages) if pages else None
            except ImportError:
                return None
        elif ext in (".docx", ".doc"):
            try:
                from docx import Document
                doc = Document(path)
                texts = [p.text for p in doc.paragraphs if p.text]
                return "\n".join(texts) if texts else None
            except ImportError:
                return None
        elif ext in (".xlsx", ".xls"):
            lines = []
            try:
                from openpyxl import load_workbook
                wb = load_workbook(path, read_only=True, data_only=True)
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    for row in ws.iter_rows(values_only=True):
                        lines.append(" | ".join(str(c) if c is not None else "" for c in row))
                return "\n".join(lines) if lines else None
            except ImportError:
                try:
                    import pandas as pd
                    dfs = pd.read_excel(path, sheet_name=None)
                    for sheet_name, df in dfs.items():
                        lines.append(f"--- Sheet: {sheet_name} ---")
                        lines.append(df.to_string())
                    return "\n".join(lines)
                except ImportError:
                    return None
        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(path)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text)
                return "\n".join(texts) if texts else None
            except ImportError:
                return None
    except Exception as e:
        logger.debug("search|knowledge|doc_extract_error|path=%s|error=%s", path, str(e)[:80])
    return None

# ────────────────────────────────────────────────────────────
# Provider Registry (9 providers)
# ────────────────────────────────────────────────────────────
SEARCH_PROVIDERS: Dict[str, Dict[str, Any]] = {
    "codecortex-codebase": {
        "id": "codecortex-codebase",
        "name": "CodeCortex Codebase Search",
        "kind": "codeSearch",
        "description": "FTS5 full-text + semantic embedding + graph-enriched symbol search across 22 languages.",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path", "repo_id",
                    "search_type", "symbol_type", "language", "file_pattern"],
    },
    "codecortex-repowt": {
        "id": "codecortex-repowt",
        "name": "CodeCortex Repository Working Tree Search",
        "kind": "repoSearch",
        "description": "Git working-tree search: staged/unstaged changes, commit history, diffs, file status.",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path", "repo_id",
                    "status_filter", "commit_range", "diff_search", "since"],
    },
    "codecortex-filesystem": {
        "id": "codecortex-filesystem",
        "name": "CodeCortex Filesystem Search",
        "kind": "fileSearch",
        "description": "Filename glob + content regex filesystem search with ReDoS protection and path validation.",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path", "repo_id",
                    "file_pattern", "content_regex", "recursive", "max_depth"],
    },
    "codecortex-graph": {
        "id": "codecortex-graph",
        "name": "CodeCortex Graph Search",
        "kind": "graphSearch",
        "description": "Graph-based symbol relationship search — callers, callees, trace flow, hierarchies.",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path", "repo_id",
                    "search_type", "max_depth", "direction", "relation_type"],
    },
    "codecortex-idegraph": {
        "id": "codecortex-idegraph",
        "name": "CodeCortex IDE Memory Search",
        "kind": "memorySearch",
        "description": "Cross-IDE conversation/memory search across 18+ IDE history formats.",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "project_name", "ide_name",
                    "search_mode", "conversation_id"],
    },
    "codecortex-knowledge": {
        "id": "codecortex-knowledge",
        "name": "CodeCortex Knowledge Graph Search",
        "kind": "knowledgeSearch",
        "description": "Documentation knowledge graph search — FTS, semantic, regex, structured queries.",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path", "repo_id",
                    "knowledge_type", "search_type"],
    },
    "codecortex-crossproject": {
        "id": "codecortex-crossproject",
        "name": "CodeCortex Cross-Project Reference Search",
        "kind": "crossProjectSearch",
        "description": "Track symbol usage across all connected repositories — functions, classes, variables.",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path", "repo_id",
                    "symbol_type", "language", "min_references"],
    },
    "codecortex-codeindex": {
        "id": "codecortex-codeindex",
        "name": "CodeCortex Code Index Data Search",
        "kind": "indexSearch",
        "description": "Fast code index metadata lookup — signatures, class definitions, structural metadata.",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path", "repo_id",
                    "symbol_type", "language", "include_signatures"],
    },
    "codecortex-agentart": {
        "id": "codecortex-agentart",
        "name": "CodeCortex .Agent Artifact Search",
        "kind": "artifactSearch",
        "description": "Search .agents folder artifacts — Coddy Codeworks versions, change history, metadata.",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path", "repo_id",
                    "artifact_type", "version", "include_history"],
    },
    "codecortex-codelogs": {
        "id": "codecortex-codelogs",
        "name": "CodeCortex Log Search",
        "kind": "logSearch",
        "description": "Search log files in <project>/logs and <project>/outputs/logs with filters for log level, time range, and free text.",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path", "repo_id",
                    "log_levels", "date_from", "date_to", "file_pattern"],
    },
    "codecortex-todo": {
        "id": "codecortex-todo",
        "name": "CodeCortex Comment Tag Scanner",
        "kind": "todoSearch",
        "description": "Detect TODO, FIXME, HACK, XXX, STUB, BUG, OPTIMIZE, REVIEW, DEPRECATED, and other comment tags across the codebase.",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path", "file_pattern", "severity"],
    },
    "codecortex-stub": {
        "id": "codecortex-stub",
        "name": "CodeCortex Stub Detector",
        "kind": "stubSearch",
        "description": "Detect empty functions, empty classes, and placeholder implementations (pass, raise NotImplementedError).",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path", "file_pattern", "language"],
    },
    "codecortex-security": {
        "id": "codecortex-security",
        "name": "CodeCortex Security Scanner",
        "kind": "securitySearch",
        "description": "Security audit: secrets (AWS keys, tokens, private keys), vulnerabilities (SQL injection, eval, pickle), PII (email, SSN), and misconfigurations (debug, CORS).",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path", "file_pattern", "severity"],
    },
    "codecortex-empty": {
        "id": "codecortex-empty",
        "name": "CodeCortex Empty File/Folder Detector",
        "kind": "emptySearch",
        "description": "Detect empty files (0 bytes), empty folders, and empty directories across the project structure.",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path", "max_depth"],
    },
    "codecortex-svn": {
        "id": "codecortex-svn",
        "name": "CodeCortex SVN Working Tree Search",
        "kind": "svnSearch",
        "description": "SVN working tree status, log history, repository info, and change tracking.",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path"],
    },
    "codecortex-blame": {
        "id": "codecortex-blame",
        "name": "CodeCortex Git Blame & Hotspot Analyzer",
        "kind": "blameSearch",
        "description": "Git blame annotations, commit hotspots, change frequency, and git history insights for every file.",
        "owned_by": "codecortex",
        "params": ["query", "max_results", "repo_path", "repo_id", "since", "commit_range"],
    },
    "codecortex-combo": {
        "id": "codecortex-combo",
        "name": "CodeCortex Unified Combo Search (17 providers)",
        "kind": "comboSearch",
        "description": "Orchestrate ALL 17 providers in parallel with graceful degradation.",
        "owned_by": "combo",
        "params": ["query", "max_results", "repo_path", "repo_id", "search_type"],
    },
}


# ────────────────────────────────────────────────────────────
# Result DTOs (9Router-compatible)
# ────────────────────────────────────────────────────────────
@dataclass
class SearchResultItem:
    title: str
    url: str
    display_url: str = ""
    snippet: str = ""
    position: int = 1
    score: float = 1.0
    published_at: Optional[str] = None
    favicon_url: Optional[str] = None
    content: Optional[str] = None
    metadata: Dict[str, Any] = field(default_factory=dict)
    citation: Dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "title": self.title,
            "url": self.url,
            "display_url": self.display_url or self.url,
            "snippet": self.snippet,
            "position": self.position,
            "score": self.score,
            "published_at": self.published_at,
            "favicon_url": self.favicon_url,
            "content": self.content,
            "metadata": self.metadata,
            "citation": self.citation,
        }


@dataclass
class SearchResponse:
    provider: str
    query: str
    results: List[SearchResultItem] = field(default_factory=list)
    answer: Optional[str] = None
    usage: Dict[str, Any] = field(default_factory=dict)
    metrics: Dict[str, Any] = field(default_factory=dict)
    errors: List[Dict[str, Any]] = field(default_factory=list)
    pagination: Optional[Dict[str, Any]] = None

    def to_dict(self) -> Dict[str, Any]:
        d = {
            "provider": self.provider,
            "query": self.query,
            "results": [r.to_dict() for r in self.results],
            "answer": self.answer,
            "usage": self.usage,
            "metrics": self.metrics,
            "errors": self.errors,
        }
        if self.pagination:
            d["pagination"] = self.pagination
        return d


@dataclass
class SearchRequest:
    query: str
    model: str = "codecortex-combo"
    max_results: int = 20
    search_type: str = "all"
    repo_path: Optional[str] = None
    repo_id: Optional[str] = None
    offset: int = 0
    # Auto-indexing flags
    auto_index: bool = True
    force_update: bool = False
    regraph: bool = False
    reindex: bool = False
    # Per-provider params
    symbol_type: str = "any"
    language: Optional[str] = None
    file_pattern: str = "*"
    content_regex: Optional[str] = None
    recursive: bool = True
    max_depth: int = 20
    search_mode: str = "keyword"
    project_name: Optional[str] = None
    ide_name: Optional[str] = None
    knowledge_type: Optional[str] = None
    direction: str = "both"
    relation_type: Optional[str] = None
    graph_max_depth: int = 3
    status_filter: Optional[str] = None
    commit_range: Optional[str] = None
    diff_search: bool = False
    since: Optional[str] = None
    min_references: int = 1
    include_signatures: bool = True
    result_filter: Optional[Dict[str, Any]] = None
    artifact_type: Optional[str] = None
    version: Optional[str] = None
    include_history: bool = False
    # Codelogs-specific params
    log_levels: Optional[str] = None
    date_from: Optional[str] = None
    date_to: Optional[str] = None


# ────────────────────────────────────────────────────────────
# Unified Search Engine
# ────────────────────────────────────────────────────────────
class UnifiedSearchEngine:
    """Orchestrate all 9 CodeCortex search providers with 9Router-compatible API."""

    def __init__(self, orchestrator: Any = None, db: Any = None):
        self._orchestrator = orchestrator
        self._db = db
        self._root_cache: Dict[str, Optional[str]] = {}

    @property
    def db(self):
        if self._db is not None:
            return self._db
        from ..core.database import DatabaseManager
        db_path = os.path.join(
            os.path.expanduser(os.getenv("CODECORTEX_DATA_DIR", os.path.join("~", ".coddy", "codecortex"))),
            "codecortex.db",
        )
        os.makedirs(os.path.dirname(db_path), exist_ok=True)
        self._db = DatabaseManager(db_path)
        # Ensure ALL tables exist (migration fallback)
        try:
            from ..core.database.migration import full_migration
            full_migration(self._db.conn)
        except Exception as e:
            logger.warning("DB migration could not run: %s", str(e)[:120])
        return self._db

    def _build_citation(self, provider: str, rank: int) -> Dict[str, Any]:
        return {
            "provider": provider,
            "retrieved_at": datetime.now(timezone.utc).isoformat(),
            "rank": rank,
        }

    def _git_root(self, path: str) -> Optional[str]:
        if path in self._root_cache:
            return self._root_cache[path]
        try:
            r = subprocess.run(
                ["git", "-C", path, "rev-parse", "--show-toplevel"],
                capture_output=True, text=True, timeout=10,
            )
            if r.returncode == 0 and r.stdout.strip():
                root = r.stdout.strip()
                self._root_cache[path] = root
                return root
        except Exception as e:
            logger.debug("search|git_root|error|path=%s|error=%s", path, str(e)[:80])
        self._root_cache[path] = None
        return None

    def _apply_result_filter(self, results: List[SearchResultItem],
                              filt: Optional[Dict[str, Any]]) -> List[SearchResultItem]:
        if not filt:
            return results
        out = []
        for r in results:
            if "source_type" in filt and r.metadata.get("source_type") not in filt["source_type"]:
                continue
            if "type" in filt and r.metadata.get("type") not in filt["type"]:
                continue
            if "language" in filt and r.metadata.get("language") != filt["language"]:
                continue
            if "min_score" in filt and r.score < filt["min_score"]:
                continue
            out.append(r)
        return out

    # ── Provider 1: codecortex-codebase ────────────────────
    async def _search_codebase(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        try:
            from ..modules.codeanalysis.services.search import Search, SearchRequest as CBRequest
            service = Search(db=self.db)
            cb_req = CBRequest(
                query=req.query, repo_id=req.repo_id, limit=req.max_results,
                file_pattern=req.file_pattern, include_content=True,
            )
            raw = service.search(cb_req, semantic=True, graph=True)
            results = []
            position = 1
            matches = raw.get("matches") or []
            for item in (matches if isinstance(matches, list) else []):
                if not isinstance(item, dict):
                    continue
                name = item.get("symbol") or str(item)
                kind = item.get("kind") or "symbol"
                fpath = item.get("file") or ""
                line = item.get("line") or 0
                sig = item.get("signature") or ""
                repo_part = f"{req.repo_id}/" if req.repo_id else ""
                results.append(SearchResultItem(
                    title=f"{kind}: {name}",
                    url=f"codecortex://repos/{repo_part}symbols/{name}",
                    display_url=f"{fpath}:{line}" if fpath else name,
                    snippet=str(sig)[:500],
                    position=position, score=float(item.get("score") or 0.9),
                    content=str(sig) if sig else None,
                    metadata={"type": kind, "language": item.get("language"),
                              "source_type": "codebase", "repo_id": req.repo_id,
                              "file_path": fpath, "line": line},
                    citation=self._build_citation("codecortex-codebase", position),
                ))
                position += 1
            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|codebase|query=%s|results=%d|elapsed=%dms",
                         req.query[:80], len(results), elapsed)
            return results, {"results_found": len(results), "total": raw.get("total_matches", len(results))}, None
        except Exception as e:
            logger.warning("search|codebase|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Provider 2: codecortex-repowt (Git Working Tree, security-filtered) ─
    async def _search_repowt(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        try:
            root = req.repo_path or os.getcwd()
            git_root = self._git_root(root)
            if not git_root:
                return [], {}, f"No git repo at: {root}"

            from ..services.security_filter import SecurityFilter
            sec_filter = SecurityFilter(project_root=git_root)
            sec_filter.load_ignore_files()

            results = []
            position = 1
            query_lower = req.query.lower() if req.query else ""

            # 1. Git status (staged + unstaged + untracked) with security filter
            try:
                sr = subprocess.run(
                    ["git", "-C", git_root, "status", "--porcelain"],
                    capture_output=True, text=True, timeout=15,
                )
                if sr.returncode == 0:
                    for line in sr.stdout.strip().split("\n"):
                        if not line.strip():
                            continue
                        status_code = line[:2].strip()
                        fname = line[3:].strip()
                        if query_lower and query_lower not in fname.lower():
                            continue

                        # SECURITY CHECK: sensitive file + ignore rules
                        fpath = os.path.join(git_root, fname)
                        sec_check = sec_filter.check_file(fpath, rel_path=fname)
                        if not sec_check["allowed"]:
                            logger.debug("search|repowt|blocked|file=%s|reasons=%s",
                                         fname, sec_check["reasons"])
                            continue

                        if status_code == "??":
                            kind = "untracked"
                        elif status_code == "!!":
                            kind = "ignored"
                        elif status_code.endswith("M"):
                            kind = "modified"
                        elif status_code.endswith("A"):
                            kind = "added"
                        elif status_code.endswith("D"):
                            kind = "deleted"
                        elif status_code.endswith("R"):
                            kind = "renamed"
                        elif status_code.endswith("C"):
                            kind = "copied"
                        else:
                            kind = "changed"
                        results.append(SearchResultItem(
                            title=f"[{status_code.strip()}] {fname}",
                            url=f"git://{git_root}/{fname}",
                            display_url=fname,
                            snippet=f"Git status: {kind}",
                            position=position, score=0.95,
                            metadata={"type": "git_status", "source_type": "repowt",
                                      "status": status_code.strip(), "kind": kind,
                                      "repo_root": git_root, "file": fname},
                            citation=self._build_citation("codecortex-repowt", position),
                        ))
                        position += 1
                        if position > req.max_results:
                            break
            except Exception as git_err:
                logger.debug("search|repowt|git_status|error=%s", str(git_err)[:80])

            # 2. Git log (commit history search)
            if position <= req.max_results and query_lower:
                log_cmd = ["git", "-C", git_root, "log", "--oneline", "--all",
                           f"--grep={req.query}", f"-{req.max_results}"]
                if req.since:
                    log_cmd.insert(3, f"--since={req.since}")
                lr = subprocess.run(log_cmd, capture_output=True, text=True, timeout=15)
                if lr.returncode == 0 and lr.stdout.strip():
                    for line in lr.stdout.strip().split("\n")[:req.max_results - position + 1]:
                        if not line.strip():
                            continue
                        parts = line.split(" ", 1)
                        commit_hash = parts[0] if parts else ""
                        commit_msg = parts[1] if len(parts) > 1 else ""
                        results.append(SearchResultItem(
                            title=f"commit {commit_hash[:8]}: {commit_msg[:100]}",
                            url=f"git://{git_root}/commit/{commit_hash}",
                            display_url=f"git:{commit_hash[:8]}",
                            snippet=commit_msg[:500],
                            position=position, score=0.8,
                            metadata={"type": "git_commit", "source_type": "repowt",
                                      "commit_hash": commit_hash, "repo_root": git_root},
                            citation=self._build_citation("codecortex-repowt", position),
                        ))
                        position += 1
            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|repowt|query=%s|root=%s|results=%d|elapsed=%dms",
                         req.query[:80], git_root, len(results), elapsed)
            return results, {"results_found": len(results), "repo_root": git_root}, None
        except Exception as e:
            logger.warning("search|repowt|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Provider 3: codecortex-filesystem (ReDoS-safe, folder-aware, security-filtered) ─
    async def _search_filesystem(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        try:
            root_path = req.repo_path or os.getcwd()
            _validate_path(root_path)
            if not os.path.isdir(root_path):
                return [], {}, f"Not a directory: {root_path}"

            # Initialize SecurityFilter with .gitignore/.aiignore
            from ..services.security_filter import SecurityFilter
            sec_filter = SecurityFilter(project_root=root_path)
            sec_filter.load_ignore_files()

            # Build default-skip set (from config) + ignored directories
            skip_dirs = {".git", "__pycache__", "node_modules", ".venv", "venv", "dist",
                         "build", ".next", ".vscode", ".idea", "target", ".tox", ".eggs"}

            results = []
            position = 1
            content_re = None
            query_lower = req.query.lower() if req.query else ""
            if req.query:
                try:
                    content_re = _safe_regex_compile(re.escape(req.query) if not req.content_regex else req.query)
                except ValueError:
                    content_re = _safe_regex_compile(re.escape(req.query))
            if req.content_regex:
                try:
                    content_re = _safe_regex_compile(req.content_regex)
                except ValueError as e:
                    return [], {}, str(e)

            file_pats = [p.strip() for p in req.file_pattern.split(",")] if req.file_pattern else ["*"]

            # Collect matching directories first (folder name matches)
            if query_lower:
                dir_matches = []
                for r, dirs, _ in os.walk(root_path):
                    depth = r.replace(root_path, "").count(os.sep)
                    if depth > req.max_depth:
                        dirs.clear()
                        continue
                    dirs[:] = [d for d in dirs if d not in skip_dirs]
                    # Apply .gitignore/.aiignore to directories
                    dirs[:] = [d for d in dirs if not sec_filter.is_ignored(os.path.relpath(os.path.join(r, d), root_path), is_dir=True)]
                    for dname in dirs:
                        if position > req.max_results:
                            break
                        if query_lower in dname.lower():
                            dpath = os.path.join(r, dname)
                            rel_path = os.path.relpath(dpath, root_path)
                            # Count files inside for snippet
                            try:
                                file_count = sum(1 for _ in Path(dpath).rglob("*") if _.is_file())
                            except Exception:
                                file_count = 0
                            dir_matches.append((dpath, rel_path, file_count))
                            break
                    if position > req.max_results:
                        break
                for dpath, rel_path, file_count in dir_matches[:req.max_results]:
                    results.append(SearchResultItem(
                        title=f"[dir] {rel_path}",
                        url=f"file://{dpath}",
                        display_url=rel_path,
                        snippet=f"Directory: {rel_path} ({file_count} files)",
                        position=position, score=0.8,
                        metadata={"type": "directory", "source_type": "filesystem",
                                  "size_bytes": 0,
                                  "modified": datetime.fromtimestamp(os.path.getmtime(dpath)).isoformat(),
                                  "file_count": file_count},
                        citation=self._build_citation("codecortex-filesystem", position),
                    ))
                    position += 1

            # Walk files (with security filtering)
            for r, dirs, files in os.walk(root_path):
                depth = r.replace(root_path, "").count(os.sep)
                if depth > req.max_depth:
                    dirs.clear()
                    continue
                dirs[:] = [d for d in dirs if d not in skip_dirs and not d.startswith(".")]
                # Apply .gitignore/.aiignore to directories
                dirs[:] = [d for d in dirs if not sec_filter.is_ignored(os.path.relpath(os.path.join(r, d), root_path), is_dir=True)]
                for fname in files:
                    if position > req.max_results:
                        break
                    if not any(fnmatch.fnmatch(fname, p) for p in file_pats):
                        continue
                    fpath = os.path.join(r, fname)
                    rel_path = os.path.relpath(fpath, root_path)

                    # SECURITY CHECK 1: Path allowed + not sensitive file + not ignored
                    sec_check = sec_filter.check_file(fpath, rel_path=rel_path)
                    if not sec_check["allowed"]:
                        logger.debug("search|filesystem|blocked|path=%s|reasons=%s",
                                     rel_path, sec_check["reasons"])
                        continue

                    # Check if filename matches query (fast path)
                    fname_match = query_lower and query_lower in fname.lower()

                    try:
                        fsize = os.path.getsize(fpath)
                        if fsize > 10 * 1024 * 1024:
                            continue
                        with open(fpath, "r", encoding="utf-8", errors="replace") as f:
                            text = f.read()
                    except Exception:
                        # Non-text file — still match if filename matches
                        if fname_match:
                            mtime = os.path.getmtime(fpath)
                            results.append(SearchResultItem(
                                title=rel_path, url=f"file://{fpath}", display_url=rel_path,
                                snippet=f"Binary/text file: {fsize} bytes",
                                position=position, score=0.6, content=None,
                                metadata={"type": "file", "source_type": "filesystem",
                                          "size_bytes": fsize,
                                          "modified": datetime.fromtimestamp(mtime).isoformat()},
                                citation=self._build_citation("codecortex-filesystem", position),
                            ))
                            position += 1
                        continue

                    # SECURITY CHECK 2: Content processing (mask/block per config)
                    processed = sec_filter.process_content(text)
                    if processed["action"] == "block":
                        logger.debug("search|filesystem|content_blocked|path=%s|reasons=%s",
                                     rel_path, processed["reasons"])
                        continue
                    elif processed["action"] == "mask":
                        # Content was masked — use masked text
                        text = processed["text"]
                        logger.debug("search|filesystem|content_masked|path=%s|reasons=%s",
                                     rel_path, processed["reasons"])

                    score = 1.0
                    snippet = ""
                    content = None
                    content_match = False
                    if content_re:
                        matches = _safe_regex_search(content_re, text)
                        if matches:
                            content_match = True
                            score = 0.7
                            match_snippets = []
                            for m in matches[:3]:
                                s = max(0, m.start() - 80)
                                e = min(len(text), m.end() + 80)
                                match_snippets.append(text[s:e].strip())
                            snippet = " | ".join(match_snippets)[:500]
                            content = text[:2000] if len(text) <= 2000 else text[:2000] + "..."
                    elif query_lower:
                        content_match = query_lower in text[:5000].lower()
                        if content_match:
                            snippet = text[:500]
                            content = text[:2000] if len(text) <= 2000 else text[:2000] + "..."
                    else:
                        snippet = text[:500]
                        content = text[:2000] if len(text) <= 2000 else text[:2000] + "..."

                    # Accept if filename matches OR content matches
                    if not fname_match and not content_match and req.query:
                        continue

                    mtime = os.path.getmtime(fpath)
                    results.append(SearchResultItem(
                        title=rel_path, url=f"file://{fpath}", display_url=rel_path,
                        snippet=snippet, position=position, score=score if content_match else 0.6,
                        content=content,
                        metadata={"type": "file", "source_type": "filesystem",
                                  "size_bytes": os.path.getsize(fpath),
                                  "modified": datetime.fromtimestamp(mtime).isoformat()},
                        citation=self._build_citation("codecortex-filesystem", position),
                    ))
                    position += 1
            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|filesystem|query=%s|dir=%s|results=%d|elapsed=%dms",
                         req.query[:80], root_path, len(results), elapsed)
            return results, {"results_found": len(results), "scanned_path": root_path}, None
        except ValueError as e:
            return [], {}, str(e)
        except Exception as e:
            logger.warning("search|filesystem|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Provider 4: codecortex-graph (auto-repo-id + auto-regraph) ─
    async def _search_graph(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        try:
            from ..core.graph import GraphManager
            from ..modules.codegraph.services.search import CODDYGraphSearch

            # Auto-detect repo_id from working tree if not provided
            effective_repo_id = req.repo_id
            if not effective_repo_id and req.repo_path:
                git_root = self._git_root(req.repo_path)
                if git_root:
                    try:
                        row = self.db.conn.execute(
                            "SELECT id FROM repositories WHERE root_path = ?",
                            (git_root,),
                        ).fetchone()
                        if row:
                            effective_repo_id = row[0] if isinstance(row, dict) else row[0]
                        else:
                            logger.debug("search|graph|repo_not_registered|git_root=%s", git_root)
                    except Exception as exc:
                        logger.warning("search|graph|repo_id_lookup_error|error=%s", str(exc)[:80])

            graph = GraphManager()
            searcher = CODDYGraphSearch(db=self.db, graph_manager=graph)

            if effective_repo_id:
                try:
                    result = await searcher.search(
                        repo_id=effective_repo_id, action="search_symbols",
                        query=req.query, limit=req.max_results,
                    )
                except Exception as e:
                    emsg = str(e)
                    if "Graph not built" in emsg or "graph" in emsg.lower() and "not built" in emsg.lower():
                        # Auto-regraph if flag set
                        if req.regraph:
                            try:
                                if hasattr(self, '_orchestrator') and self._orchestrator is not None:
                                    if hasattr(self._orchestrator, 'graph_service'):
                                        logger.info("search|graph|auto_regraph|repo=%s", effective_repo_id)
                                        await self._orchestrator.graph_service.build_graph(
                                            repo_id=effective_repo_id,
                                            detect_modular=True,
                                            build_dependency_graph=True,
                                        )
                                        # Re-try search after regraph
                                        result = await searcher.search(
                                            repo_id=effective_repo_id, action="search_symbols",
                                            query=req.query, limit=req.max_results,
                                        )
                                    else:
                                        logger.warning("search|graph|no_graph_service|orchestrator missing graph_service")
                                        return [], {"results_found": 0, "warning": f"Graph service not available for repo '{effective_repo_id}'"}, None
                                else:
                                    # Try importing orchestrator
                                    try:
                                        from ..main import CortexOrchestrator
                                        orch = CortexOrchestrator()
                                        logger.info("search|graph|auto_regraph_lazy|repo=%s", effective_repo_id)
                                        await orch.graph_service.build_graph(
                                            repo_id=effective_repo_id,
                                            detect_modular=True,
                                            build_dependency_graph=True,
                                        )
                                        result = await searcher.search(
                                            repo_id=effective_repo_id, action="search_symbols",
                                            query=req.query, limit=req.max_results,
                                        )
                                    except Exception as ge:
                                        logger.warning("search|graph|auto_regraph_failed|repo=%s|error=%s", effective_repo_id, str(ge)[:80])
                                        return [], {"results_found": 0, "warning": f"Graph build failed for repo '{effective_repo_id}': {str(ge)[:80]}. Try manual: codecortex cg build"}, None
                            except Exception as ge:
                                logger.warning("search|graph|auto_regraph_error|repo=%s|error=%s", effective_repo_id, str(ge)[:80])
                                return [], {"results_found": 0, "warning": f"Graph needs building for repo '{effective_repo_id}'. Run: codecortex cg build"}, None
                        else:
                            return [], {"results_found": 0, "warning": f"Graph needs building for repo '{effective_repo_id}'. Set regraph=True or run: codecortex cg build"}, None
                    return [], {}, emsg
            else:
                if req.repo_path:
                    hint = f"Repository at '{req.repo_path}' not registered. Run: codecortex repo analyze {req.repo_path}"
                else:
                    hint = "Provide --repo-id or --repo-path to enable graph search"
                return [], {"results_found": 0, "warning": hint}, None

            results = []
            position = 1
            items = result.get("results") if isinstance(result, dict) else (
                getattr(result, "results", []) if hasattr(result, "results") else [])
            for item in (items if isinstance(items, list) else [])[:req.max_results]:
                if isinstance(item, dict):
                    name = item.get("name") or str(item)
                    kind = item.get("kind") or "node"
                    fpath = item.get("file") or ""
                    line = item.get("line") or 0
                else:
                    name, kind, fpath, line = str(item), "node", "", 0
                repo_part = f"{effective_repo_id}/" if effective_repo_id else ""
                results.append(SearchResultItem(
                    title=f"{kind}: {name}",
                    url=f"codecortex://repos/{repo_part}graph/{kind}/{name}",
                    display_url=f"{fpath}:{line}" if fpath else name,
                    snippet=str(item)[:500], position=position, score=0.85,
                    metadata={"type": kind, "source_type": "graph", "repo_id": effective_repo_id},
                    citation=self._build_citation("codecortex-graph", position),
                ))
                position += 1
            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|graph|query=%s|repo=%s|results=%d|elapsed=%dms",
                         req.query[:80], effective_repo_id, len(results), elapsed)
            return results, {"results_found": len(results)}, None
        except Exception as e:
            logger.warning("search|graph|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Provider 5: codecortex-idegraph (with auto-harvest) ─
    async def _search_idegraph(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        try:
            from ..modules.idegraph.services.search import Search as IGSearch
            svc = IGSearch(db=self.db)
            engrams = svc.search(query=req.query, project_name=req.project_name,
                                ide_name=req.ide_name, limit=req.max_results)
            results = []
            position = 1
            for eng in engrams[:req.max_results]:
                cid = getattr(eng, 'id', '') or str(eng)
                title = getattr(eng, 'title', '') or 'Memory'
                ide = getattr(eng.ide_info, 'name', '') if hasattr(eng, 'ide_info') and eng.ide_info else 'unknown'
                snippet = ''
                if hasattr(eng, 'messages') and eng.messages:
                    snippet = (getattr(eng.messages[0], 'content', '') or '')[:500]
                results.append(SearchResultItem(
                    title=str(title)[:200], url=f"codecortex://idegraph/memories/{cid}",
                    display_url=f"ide:{ide}", snippet=str(snippet)[:500],
                    position=position, score=0.7, content=str(snippet),
                    metadata={"type": "ide_memory", "source_type": "idegraph",
                              "ide": ide, "project": getattr(eng, 'project_name', None),
                              "conversation_id": cid},
                    citation=self._build_citation("codecortex-idegraph", position),
                ))
                position += 1

            # Auto-harvest: if no results found, trigger harvest then re-run
            if len(results) == 0 and req.repo_path and os.path.isdir(req.repo_path):
                logger.info("search|idegraph|auto_harvest|path=%s", req.repo_path)
                try:
                    from ..modules.idegraph.core.orchestrator import SideCortexOrchestrator
                    orch = SideCortexOrchestrator()
                    engram_results = await asyncio.to_thread(orch.run_all)
                    logger.info("search|idegraph|auto_harvest|complete|path=%s|engrams=%d", req.repo_path, len(engram_results))
                    # Re-run search after harvest
                    engrams = svc.search(query=req.query, project_name=req.project_name,
                                        ide_name=req.ide_name, limit=req.max_results)
                    for eng in engrams[:req.max_results]:
                        cid = getattr(eng, 'id', '') or str(eng)
                        title = getattr(eng, 'title', '') or 'Memory'
                        ide = getattr(eng.ide_info, 'name', '') if hasattr(eng, 'ide_info') and eng.ide_info else 'unknown'
                        snippet = ''
                        if hasattr(eng, 'messages') and eng.messages:
                            snippet = (getattr(eng.messages[0], 'content', '') or '')[:500]
                        results.append(SearchResultItem(
                            title=str(title)[:200], url=f"codecortex://idegraph/memories/{cid}",
                            display_url=f"ide:{ide}", snippet=str(snippet)[:500],
                            position=position, score=0.7, content=str(snippet),
                            metadata={"type": "ide_memory", "source_type": "idegraph",
                                      "ide": ide, "project": getattr(eng, 'project_name', None),
                                      "conversation_id": cid},
                            citation=self._build_citation("codecortex-idegraph", position),
                        ))
                        position += 1
                except Exception as e:
                    logger.warning("search|idegraph|auto_harvest|error=%s", str(e)[:80])

            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|idegraph|query=%s|results=%d|elapsed=%dms",
                         req.query[:80], len(results), elapsed)
            return results, {"results_found": len(results)}, None
        except Exception as e:
            logger.warning("search|idegraph|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Provider 6: codecortex-knowledge (DB + Markdown fallback) ─
    async def _search_knowledge(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        results = []
        position = 1
        try:
            from ..modules.knowledgegraph.adapters.storage import KnowledgeStore
            store = KnowledgeStore(db=self.db)
            k_results = store.query(
                fts_query=req.query, limit=req.max_results,
                knowledge_types=[req.knowledge_type] if req.knowledge_type else None,
            )
            items = k_results.get("chunks") if isinstance(k_results, dict) else []
            for item in (items if isinstance(items, list) else [])[:req.max_results]:
                if isinstance(item, dict):
                    kid = item.get("id") or str(item)
                    title = item.get("title") or "Knowledge Item"
                    content = item.get("content") or ""
                else:
                    kid = getattr(item, 'id', str(item))
                    title = getattr(item, 'title', '') or 'Knowledge'
                    content = getattr(item, 'content', '') or ''
                snippet = str(content)[:500]
                results.append(SearchResultItem(
                    title=str(title)[:200], url=f"codecortex://knowledge/items/{kid}",
                    display_url=str(title)[:100], snippet=snippet,
                    position=position, score=0.6, content=snippet,
                    metadata={"type": "knowledge_chunk", "source_type": "knowledge",
                              "id": str(kid)},
                    citation=self._build_citation("codecortex-knowledge", position),
                ))
                position += 1
        except Exception as e:
            logger.debug("knowledge DB query failed, falling back to doc scan: %s", str(e)[:80])

        # Fallback: search all document types (text + binary) with security filtering
        if len(results) < req.max_results:
            try:
                root = req.repo_path or os.getcwd()
                if os.path.isdir(root):
                    from ..services.security_filter import SecurityFilter
                    sec_filter = SecurityFilter(project_root=root)
                    sec_filter.load_ignore_files()
                    text_exts = {".md", ".rst", ".txt", ".markdown", ".adoc", ".org", ".csv", ".log"}
                    binary_exts = {".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx"}
                    ALLOWED_BY_NAME = {"readme", "changelog", "license", "contributing"}
                    query_lower = req.query.lower()
                    skip_dirs = {".git", "__pycache__", "node_modules", ".venv", "venv", "dist",
                                 "build", ".next", ".vscode", ".idea", "target", ".tox", ".eggs"}
                    for r, dirs, files in os.walk(root):
                        depth = r.replace(root, "").count(os.sep)
                        if depth > 10:
                            dirs.clear()
                            continue
                        dirs[:] = [d for d in dirs if not d.startswith(".") and d not in skip_dirs]
                        # Apply .gitignore/.aiignore to directories
                        dirs[:] = [d for d in dirs if not sec_filter.is_ignored(os.path.relpath(os.path.join(r, d), root), is_dir=True)]
                        for fname in files:
                            if position > req.max_results:
                                break
                            ext = os.path.splitext(fname)[1].lower()
                            is_text = ext in text_exts
                            is_binary = ext in binary_exts
                            if not is_text and not is_binary and fname.lower().replace("_", "").replace("-", "").replace(".md", "") not in ALLOWED_BY_NAME:
                                continue
                            fpath = os.path.join(r, fname)
                            rel_path = os.path.relpath(fpath, root)

                            # SECURITY CHECK: Path allowed + not sensitive file + not ignored
                            sec_check = sec_filter.check_file(fpath, rel_path=rel_path)
                            if not sec_check["allowed"]:
                                logger.debug("search|knowledge|blocked|path=%s|reasons=%s",
                                             rel_path, sec_check["reasons"])
                                continue

                            try:
                                if os.path.getsize(fpath) > 5 * 1024 * 1024:
                                    continue
                                if is_text:
                                    with open(fpath, "r", encoding="utf-8", errors="replace") as f:
                                        text = f.read()
                                elif is_binary:
                                    text = _extract_document_text(str(fpath)) or ""
                                else:
                                    try:
                                        with open(fpath, "r", encoding="utf-8", errors="replace") as f:
                                            text = f.read()
                                    except Exception:
                                        text = _extract_document_text(str(fpath)) or ""
                            except Exception:
                                if is_binary:
                                    text = _extract_document_text(str(fpath)) or ""
                                else:
                                    continue
                            if not text:
                                continue

                            # SECURITY CHECK: Content processing (mask/block per config)
                            processed = sec_filter.process_content(text)
                            if processed["action"] == "block":
                                logger.debug("search|knowledge|content_blocked|path=%s|reasons=%s",
                                             rel_path, processed["reasons"])
                                continue
                            elif processed["action"] == "mask":
                                text = processed["text"]
                                logger.debug("search|knowledge|content_masked|path=%s|reasons=%s",
                                             rel_path, processed["reasons"])

                            if query_lower and query_lower not in text[:10000].lower() and query_lower not in fname.lower():
                                continue
                            rel_path = os.path.relpath(fpath, root)
                            idx = text.lower().find(query_lower) if query_lower else 0
                            snippet_start = max(0, idx - 200)
                            snippet = text[snippet_start:snippet_start + 800]
                            results.append(SearchResultItem(
                                title=f"[{ext[1:] if ext else 'doc'}] {rel_path}",
                                url=f"file://{fpath}",
                                display_url=rel_path,
                                snippet=snippet[:500],
                                position=position, score=0.55,
                                content=text[:2000] if len(text) <= 2000 else text[:2000] + "...",
                                metadata={"type": "documentation", "source_type": "knowledge",
                                          "file_path": rel_path, "format": ext,
                                          "size_bytes": os.path.getsize(fpath),
                                          "modified": datetime.fromtimestamp(os.path.getmtime(fpath)).isoformat()},
                                citation=self._build_citation("codecortex-knowledge", position),
                            ))
                            position += 1
            except Exception as e:
                logger.debug("knowledge markdown fallback failed: %s", str(e)[:80])

        elapsed = int((time.monotonic() - t0) * 1000)
        logger.info("search|knowledge|query=%s|results=%d|elapsed=%dms",
                     req.query[:80], len(results), elapsed)
        return results, {"results_found": len(results)}, None

    # ── Provider 7: codecortex-crossproject ────────────────
    async def _search_crossproject(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        try:
            db_conn = self.db.conn
            all_repos = db_conn.execute(
                "SELECT id, name, root_path FROM repositories"
            ).fetchall() if hasattr(db_conn, 'execute') else []
            results = []
            position = 1
            query_lower = req.query.lower()
            symbol_type = req.symbol_type if req.symbol_type != "any" else None

            # Check if data exists for warning
            try:
                total_symbols = db_conn.execute("SELECT COUNT(*) FROM symbols").fetchone()
                total_symbols = total_symbols[0] if total_symbols else 0
            except Exception:
                total_symbols = 0

            if total_symbols == 0:
                msg = "No symbols indexed yet. Run: codecortex repo analyze <path> to build index"
                logger.warning("search|crossproject|warning=%s", msg)
                return [], {"results_found": 0, "warning": msg, "indexed": False}, None

            for repo_row in all_repos:
                rid = repo_row["id"] if isinstance(repo_row, dict) else repo_row[0]
                rname = repo_row.get("name", "") if isinstance(repo_row, dict) else (repo_row[1] if len(repo_row) > 1 else "")
                try:
                    rows = db_conn.execute(
                        """SELECT s.name, s.symbol_type, s.start_line, s.signature,
                                  f.relative_path, s.id
                           FROM symbols s LEFT JOIN files f ON s.file_id = f.id
                           WHERE s.repository_id = ? AND LOWER(s.name) LIKE ?
                           ORDER BY s.name LIMIT ?""",
                        (rid, f"%{query_lower}%", req.max_results),
                    ).fetchall()
                except Exception:
                    continue
                for row in rows:
                    name = row["name"] if isinstance(row, dict) else row[0]
                    stype = row["symbol_type"] if isinstance(row, dict) else row[1]
                    fpath = (row.get("relative_path") or "") if isinstance(row, dict) else (row[4] or "")
                    line = row["start_line"] if isinstance(row, dict) else row[2]
                    sig = (row.get("signature") or "") if isinstance(row, dict) else (row[3] or "")
                    if symbol_type and stype != symbol_type:
                        continue
                    if req.language and isinstance(row, dict):
                        lang = row.get("language", "")
                        if req.language.lower() not in lang.lower():
                            continue
                    results.append(SearchResultItem(
                        title=f"[{rname}] {stype}: {name}",
                        url=f"codecortex://repos/{rid}/symbols/{name}",
                        display_url=f"{rname}/{fpath}:{line}" if fpath else f"{rname}/{name}",
                        snippet=str(sig)[:500],
                        position=position, score=0.75, content=str(sig) if sig else None,
                        metadata={"type": stype, "source_type": "crossproject",
                                  "repo_id": rid, "repo_name": rname,
                                  "file_path": fpath, "line": line},
                        citation=self._build_citation("codecortex-crossproject", position),
                    ))
                    position += 1
                    if position > req.max_results:
                        break
                if position > req.max_results:
                    break
            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|crossproject|query=%s|repos=%d|results=%d|elapsed=%dms",
                         req.query[:80], len(all_repos), len(results), elapsed)
            return results, {"results_found": len(results), "repos_scanned": len(all_repos)}, None
        except Exception as e:
            logger.warning("search|crossproject|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Provider 8: codecortex-codeindex ───────────────────
    async def _search_codeindex(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        """Fast code index metadata lookup from indexed symbols table."""
        t0 = time.monotonic()
        try:
            db_conn = self.db.conn
            query_lower = f"%{req.query.lower()}%"
            symbol_type = req.symbol_type if req.symbol_type != "any" else None

            # Check indexing status
            try:
                total_symbols = db_conn.execute("SELECT COUNT(*) FROM symbols").fetchone()
                total_symbols = total_symbols[0] if total_symbols else 0
            except Exception:
                total_symbols = 0

            if total_symbols == 0:
                msg = "Code index is empty. Run: codecortex repo analyze <path> to build index"
                logger.warning("search|codeindex|warning=%s", msg)
                return [], {"results_found": 0, "warning": msg, "indexed": False}, None

            sql = """SELECT s.name, s.symbol_type, s.start_line, s.signature,
                            s.docstring, f.relative_path, s.id, s.repository_id
                     FROM symbols s LEFT JOIN files f ON s.file_id = f.id
                     WHERE LOWER(s.name) LIKE ?
                     ORDER BY s.name LIMIT ?"""
            params: list = [query_lower, req.max_results]

            if req.repo_id:
                sql = sql.replace("WHERE", "WHERE s.repository_id = ? AND")
                params = [req.repo_id, query_lower, req.max_results]

            rows = db_conn.execute(sql, params).fetchall()
            results = []
            position = 1
            for row in rows:
                name = row["name"] if isinstance(row, dict) else row[0]
                stype = row["symbol_type"] if isinstance(row, dict) else row[1]
                line = row["start_line"] if isinstance(row, dict) else row[2]
                sig = row.get("signature", "") if isinstance(row, dict) else (row[3] or "")
                fpath = row.get("relative_path", "") if isinstance(row, dict) else (row[5] or "")
                rid = row.get("repository_id", "") if isinstance(row, dict) else (row[7] or "")

                if symbol_type and stype != symbol_type:
                    continue
                if req.language and isinstance(row, dict):
                    lang = row.get("language", "")
                    if req.language.lower() not in lang.lower():
                        continue

                repo_part = f"{rid}/" if rid else ""
                results.append(SearchResultItem(
                    title=f"{stype}: {name}",
                    url=f"codecortex://repos/{repo_part}index/{stype}/{name}",
                    display_url=f"{fpath}:{line}" if fpath else name,
                    snippet=str(sig)[:500] if req.include_signatures else f"{stype} {name} at {fpath}:{line}",
                    position=position, score=0.88,
                    content=str(sig) if req.include_signatures and sig else None,
                    metadata={"type": stype, "source_type": "codeindex",
                              "repo_id": rid, "file_path": fpath, "line": line},
                    citation=self._build_citation("codecortex-codeindex", position),
                ))
                position += 1
            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|codeindex|query=%s|results=%d|elapsed=%dms",
                         req.query[:80], len(results), elapsed)
            return results, {"results_found": len(results)}, None
        except Exception as e:
            logger.warning("search|codeindex|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Provider 10: codecortex-codelogs (enhanced — systematic path discovery) ─
    async def _search_codelogs(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        try:
            from src.modules.codelogs.services.log_service import LogService, LogSearchFilter

            root = req.repo_path or os.getcwd()
            search_paths = getattr(req, 'file_pattern', None)
            svc = LogService(project_root=root)

            level_list = [l.strip().upper() for l in req.log_levels.split(",")] if req.log_levels else None
            filt = LogSearchFilter(
                query=req.query,
                log_levels=level_list,
                date_from=req.date_from,
                date_to=req.date_to,
                file_pattern=req.file_pattern,
                max_results=req.max_results,
                offset=req.offset,
            )
            entries = svc.search(filt)

            results = []
            for position, entry in enumerate(entries, 1):
                snippet = entry.message[:500]
                results.append(SearchResultItem(
                    title=f"[{entry.level}] {entry.path}:{entry.line}",
                    url=f"file://{os.path.join(root, entry.path)}",
                    display_url=f"{entry.path}:{entry.line}",
                    snippet=snippet,
                    content=entry.message[:2000] if entry.message else None,
                    position=position,
                    score=entry.score,
                    published_at=entry.timestamp,
                    metadata={
                        "type": "log_entry",
                        "source_type": "codelogs",
                        "log_level": entry.level,
                        "file_path": entry.path,
                        "line": entry.line,
                        "timestamp": entry.timestamp,
                    },
                    citation=self._build_citation("codecortex-codelogs", position),
                ))
            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|codelogs|query=%s|root=%s|results=%d|elapsed=%dms",
                         req.query[:80], root, len(results), elapsed)
            return results, {"results_found": len(results), "scanned_root": root}, None
        except Exception as e:
            logger.warning("search|codelogs|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Provider 11: codecortex-todo (comment tag scanner) ──────────
    async def _search_todo(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        try:
            from ..modules.codeanalysis.analyzers.audit import CodeAuditor

            root = req.repo_path or os.getcwd()
            if not os.path.isdir(root):
                return [], {}, f"Not a directory: {root}"

            tag_filter = req.query.strip().upper() if req.query else None
            result = CodeAuditor.audit({
                "target_path": root,
                "scan_categories": ["comments"],
                "severity_threshold": req.result_filter.get("min_severity", "low") if req.result_filter else "low",
                "max_file_size_kb": 1024,
                "use_aiignore": True,
            })
            if not isinstance(result, dict) or not result.get("success"):
                return [], {}, result.get("message", "Audit failed")

            findings = result.get("data", {}).get("findings", {}).get("comment_tags", [])
            results = []
            position = 1
            for f in findings:
                if position > req.max_results:
                    break
                tag = f.get("tag", "")
                if tag_filter and tag_filter != tag and tag_filter not in f.get("message", "").upper():
                    continue
                fpath = f.get("file", "")
                line = f.get("line", 0)
                results.append(SearchResultItem(
                    title=f"[{tag}] {fpath}:{line}",
                    url=f"file://{fpath}",
                    display_url=f"{fpath}:{line}",
                    snippet=f.get("message", "")[:500],
                    content=f.get("context", "")[:2000] if f.get("context") else None,
                    position=position, score=f.get("confidence", 0.85),
                    metadata={
                        "type": "comment_tag",
                        "source_type": "todo",
                        "tag": tag,
                        "priority": f.get("priority", "medium"),
                        "file_path": fpath, "line": line,
                        "confidence": f.get("confidence", 0.85),
                        "comment_type": f.get("comment_type", "line"),
                    },
                    citation=self._build_citation("codecortex-todo", position),
                ))
                position += 1
            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|todo|path=%s|tags=%d|results=%d|elapsed=%dms",
                         root, len(findings), len(results), elapsed)
            return results, {"results_found": len(results), "total_tags": len(findings)}, None
        except Exception as e:
            logger.warning("search|todo|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Provider 12: codecortex-stub (empty function/class detector) ─
    async def _search_stub(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        try:
            from ..modules.codeanalysis.analyzers.audit import CodeAuditor

            root = req.repo_path or os.getcwd()
            if not os.path.isdir(root):
                return [], {}, f"Not a directory: {root}"

            result = CodeAuditor.audit({
                "target_path": root,
                "scan_categories": ["empty_code"],
                "severity_threshold": "low",
                "max_file_size_kb": 1024,
                "use_aiignore": True,
            })
            if not isinstance(result, dict) or not result.get("success"):
                return [], {}, result.get("message", "Audit failed")

            data = result.get("data", {})
            empty_funcs = data.get("findings", {}).get("empty_functions", [])
            empty_classes = data.get("findings", {}).get("empty_classes", [])
            findings = empty_funcs + empty_classes
            results = []
            position = 1
            lang = (req.language or "").lower()
            for f in findings:
                if position > req.max_results:
                    break
                fpath = f.get("file", "")
                line = f.get("line", 0)
                ftype = f.get("type", "stub")
                if lang and lang not in fpath.lower():
                    continue
                results.append(SearchResultItem(
                    title=f"[{ftype}] {fpath}:{line}",
                    url=f"file://{fpath}", display_url=f"{fpath}:{line}",
                    snippet=f.get("message", "")[:500],
                    position=position, score=f.get("confidence", 0.9),
                    metadata={
                        "type": ftype,
                        "source_type": "stub",
                        "priority": f.get("severity", "medium"),
                        "file_path": fpath, "line": line,
                        "confidence": f.get("confidence", 0.9),
                    },
                    citation=self._build_citation("codecortex-stub", position),
                ))
                position += 1
            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|stub|path=%s|results=%d|elapsed=%dms",
                         root, len(results), elapsed)
            return results, {"results_found": len(results), "total_stubs": len(findings)}, None
        except Exception as e:
            logger.warning("search|stub|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Provider 13: codecortex-security (secrets, vulns, PII, misconfig) ─
    async def _search_security(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        try:
            from ..modules.codeanalysis.analyzers.audit import CodeAuditor
            from ..modules.filesystem.adapters.audit import DiskAudit

            root = req.repo_path or os.getcwd()
            if not os.path.isdir(root):
                return [], {}, f"Not a directory: {root}"

            # Phase 1: Content audit via CodeAuditor
            severity_threshold = req.result_filter.get("min_severity", "low") if req.result_filter else "low"
            audit_result = CodeAuditor.audit({
                "target_path": root,
                "scan_categories": ["secrets", "vulns", "pii", "misconfig"],
                "severity_threshold": severity_threshold,
                "max_file_size_kb": 2048,
                "use_aiignore": True,
            })
            findings = []
            if isinstance(audit_result, dict) and audit_result.get("success"):
                cats = audit_result.get("data", {}).get("findings", {})
                for cat_name in ("secrets", "vulnerabilities", "pii", "misconfig"):
                    findings.extend(cats.get(cat_name, []))

            # Phase 2: Sensitive file detection via DiskAudit
            try:
                fs_audit = DiskAudit.audit({
                    "target": root, "recursive": True,
                    "severity": [severity_threshold],
                    "check_permissions": True, "check_hidden": True,
                    "max_file_size_mb": 100, "limit": 100,
                })
                fs_findings = fs_audit.get("data", {}).get("findings", [])
                for ff in fs_findings:
                    findings.append({
                        "type": "sensitive_file",
                        "severity": ff.get("severity", "medium"),
                        "line": 0,
                        "file": ff.get("path", ""),
                        "message": ff.get("reason", "Sensitive file detected"),
                        "details": {"category": ff.get("category", ""),
                                    "recommendation": ff.get("recommendation", "")},
                        "confidence": 0.9,
                    })
            except Exception as fe:
                logger.warning("search|security|disk_audit_error=%s", str(fe)[:80])

            results = []
            position = 1
            query_lower = req.query.lower() if req.query else ""
            for f in findings:
                if position > req.max_results:
                    break
                fpath = f.get("file", "")
                line = f.get("line", 0)
                ftype = f.get("type", "security_issue")
                msg = f.get("message", "")
                if query_lower and query_lower not in msg.lower() and query_lower not in ftype.lower():
                    continue
                snippet = msg[:500]
                if isinstance(f.get("details"), dict) and f["details"].get("recommendation"):
                    snippet += f" | recommendation: {f['details']['recommendation']}"
                results.append(SearchResultItem(
                    title=f"[{ftype}] {fpath}:{line}" if line else f"[{ftype}] {fpath}",
                    url=f"file://{fpath}" if fpath else "",
                    display_url=f"{fpath}:{line}" if line else fpath,
                    snippet=snippet,
                    position=position, score=f.get("confidence", 0.85),
                    metadata={
                        "type": ftype,
                        "source_type": "security",
                        "priority": f.get("severity", "medium"),
                        "file_path": fpath, "line": line,
                        "confidence": f.get("confidence", 0.85),
                    },
                    citation=self._build_citation("codecortex-security", position),
                ))
                position += 1
            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|security|path=%s|results=%d|elapsed=%dms",
                         root, len(results), elapsed)
            return results, {"results_found": len(results), "total_findings": len(findings)}, None
        except Exception as e:
            logger.warning("search|security|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Provider 14: codecortex-empty (empty files + folders) ─────────
    async def _search_empty(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        try:
            from ..modules.codeanalysis.analyzers.audit import CodeAuditor

            root = req.repo_path or os.getcwd()
            if not os.path.isdir(root):
                return [], {}, f"Not a directory: {root}"

            skip_dirs = {".git", "__pycache__", "node_modules", ".venv", "venv",
                         "dist", "build", ".next", ".vscode", ".idea", "target",
                         ".tox", ".eggs", ".svn"}
            results = []
            position = 1
            max_depth = req.max_depth
            query_lower = req.query.lower() if req.query else ""
            empty_dirs: List[str] = []

            for r, dirs, files in os.walk(root):
                depth = r.replace(root, "").count(os.sep)
                if max_depth and depth > max_depth:
                    dirs.clear()
                    continue
                dirs[:] = [d for d in dirs if d not in skip_dirs]

                # Empty directory detection (no files inside after filtering)
                rel_dir = os.path.relpath(r, root)
                visible_files = [f for f in files if not f.startswith(".") and not f.endswith((".pyc", ".pyo"))]
                visible_dirs = [d for d in dirs if not d.startswith(".")]
                if not visible_files and not visible_dirs and rel_dir != ".":
                    empty_dirs.append((rel_dir, r))

                # Empty file detection (0 bytes)
                for fname in files:
                    if position > req.max_results:
                        break
                    if fname.startswith("."):
                        continue
                    fpath = os.path.join(r, fname)
                    try:
                        if os.path.getsize(fpath) == 0:
                            rel = os.path.relpath(fpath, root)
                            if query_lower and query_lower not in rel.lower():
                                continue
                            results.append(SearchResultItem(
                                title=f"[empty_file] {rel}",
                                url=f"file://{fpath}",
                                display_url=rel,
                                snippet="File is empty (0 bytes)",
                                position=position, score=0.7,
                                metadata={
                                    "type": "empty_file",
                                    "source_type": "empty",
                                    "file_path": rel, "line": 0,
                                    "size_bytes": 0,
                                },
                                citation=self._build_citation("codecortex-empty", position),
                            ))
                            position += 1
                    except OSError:
                        continue
                if position > req.max_results:
                    break

            # Append empty directories
            for rel_dir, abs_dir in empty_dirs:
                if position > req.max_results:
                    break
                if query_lower and query_lower not in rel_dir.lower():
                    continue
                results.append(SearchResultItem(
                    title=f"[empty_dir] {rel_dir}",
                    url=f"file://{abs_dir}",
                    display_url=rel_dir,
                    snippet="Directory contains no files",
                    position=position, score=0.6,
                    metadata={
                        "type": "empty_directory",
                        "source_type": "empty",
                        "file_path": rel_dir, "line": 0,
                        "size_bytes": 0,
                    },
                    citation=self._build_citation("codecortex-empty", position),
                ))
                position += 1

            # Also detect empty function/class stubs for completeness
            try:
                stub_result = CodeAuditor.audit({
                    "target_path": root,
                    "scan_categories": ["empty_code"],
                    "severity_threshold": "low",
                    "max_file_size_kb": 1024,
                })
                if isinstance(stub_result, dict) and stub_result.get("success"):
                    for cat in ("empty_functions", "empty_classes"):
                        for f in stub_result["data"]["findings"].get(cat, []):
                            if position > req.max_results:
                                break
                            fpath = f.get("file", "")
                            line = f.get("line", 0)
                            ftype = f.get("type", "stub")
                            rel = os.path.relpath(fpath, root) if fpath else ""
                            if query_lower and query_lower not in rel.lower():
                                continue
                            results.append(SearchResultItem(
                                title=f"[{ftype}] {rel}:{line}" if line else f"[{ftype}] {rel}",
                                url=f"file://{fpath}",
                                display_url=f"{rel}:{line}" if line else rel,
                                snippet=f.get("message", "")[:500],
                                position=position, score=f.get("confidence", 0.85),
                                metadata={
                                    "type": ftype,
                                    "source_type": "empty",
                                    "file_path": rel, "line": line,
                                    "priority": f.get("severity", "medium"),
                                },
                                citation=self._build_citation("codecortex-empty", position),
                            ))
                            position += 1
            except Exception as se:
                logger.debug("search|empty|stub_scan_skipped=%s", str(se)[:80])

            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|empty|path=%s|results=%d|empty_dirs=%d|elapsed=%dms",
                         root, len(results), len(empty_dirs), elapsed)
            return results, {"results_found": len(results), "empty_dirs_found": len(empty_dirs)}, None
        except Exception as e:
            logger.warning("search|empty|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Provider 15: codecortex-svn (Subversion working tree) ──────────
    async def _search_svn(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        try:
            from ..modules.filesystem.adapters.svn import DiskSvn

            root = req.repo_path or os.getcwd()
            root_path = Path(root).resolve()
            svn_root = DiskSvn.find_root(root_path)
            if not svn_root:
                return [], {}, f"Not an SVN working copy: {root}"
            if not DiskSvn.is_svn_available():
                return [], {}, "SVN CLI not available"

            results = []
            position = 1

            # Phase 1: Repository insights
            if position <= req.max_results:
                try:
                    insights = DiskSvn.get_insights(root_path)
                    if insights:
                        snippet_parts = []
                        for k in ("revision", "url", "last_author", "last_changed_date", "repository_root"):
                            v = insights.get(k)
                            if v:
                                snippet_parts.append(f"{k}={v}")
                        results.append(SearchResultItem(
                            title="[svn_info] Repository Overview",
                            url=insights.get("url", ""),
                            display_url=str(svn_root),
                            snippet=" | ".join(snippet_parts),
                            position=position, score=0.95,
                            metadata={
                                "type": "svn_info",
                                "source_type": "svn",
                                "file_path": str(svn_root), "line": 0,
                                "revision": insights.get("revision"),
                                "url": insights.get("url", ""),
                                "last_author": insights.get("last_author", ""),
                            },
                            citation=self._build_citation("codecortex-svn", position),
                        ))
                        position += 1
                except Exception as ie:
                    logger.debug("search|svn|insights_skipped=%s", str(ie)[:80])

            # Phase 2: Status (modified, added, deleted, etc.)
            if position <= req.max_results:
                try:
                    status_result = subprocess.run(
                        ["svn", "status", "--non-interactive"],
                        cwd=str(svn_root), capture_output=True, text=True, timeout=15,
                    )
                    if status_result.returncode == 0:
                        query_lower = req.query.lower() if req.query else ""
                        for line in status_result.stdout.splitlines():
                            if position > req.max_results:
                                break
                            if len(line) < 8:
                                continue
                            svn_code = line[0]
                            fpath = line[7:].strip()
                            if query_lower and query_lower not in fpath.lower():
                                continue
                            status_map = {
                                "M": "modified", "A": "added", "D": "deleted",
                                "R": "replaced", "C": "conflicted", "?": "unversioned",
                                "!": "missing", "~": "obstructed", "L": "locked",
                            }
                            status = status_map.get(svn_code, svn_code)
                            results.append(SearchResultItem(
                                title=f"[{status}] {fpath}",
                                url=f"file://{svn_root / fpath}",
                                display_url=fpath,
                                snippet=f"SVN status: {status}",
                                position=position, score=0.9,
                                metadata={
                                    "type": "svn_status",
                                    "source_type": "svn",
                                    "file_path": fpath, "line": 0,
                                    "svn_status": status,
                                },
                                citation=self._build_citation("codecortex-svn", position),
                            ))
                            position += 1
                except Exception as se:
                    logger.debug("search|svn|status_skipped=%s", str(se)[:80])

            # Phase 3: Log history matching query
            if req.query and position <= req.max_results:
                try:
                    log_result = subprocess.run(
                        ["svn", "log", "-l", "50", "--non-interactive"],
                        cwd=str(svn_root), capture_output=True, text=True, timeout=15,
                    )
                    if log_result.returncode == 0:
                        query_lower = req.query.lower()
                        current_rev = {}
                        for line in log_result.stdout.splitlines():
                            m = re.match(r"r(\d+)\s*\|\s*(\S+)\s*\|\s*(.*?)\s*\|\s*(\d+)", line)
                            if m:
                                if current_rev.get("message") and query_lower in current_rev["message"].lower():
                                    results.append(SearchResultItem(
                                        title=f"[svn_log] r{current_rev['revision']}",
                                        url="", display_url=f"r{current_rev['revision']} by {current_rev['author']}",
                                        snippet=current_rev["message"][:500],
                                        position=position, score=0.85,
                                        metadata={
                                            "type": "svn_log",
                                            "source_type": "svn",
                                            "file_path": "svn history", "line": 0,
                                            "revision": current_rev["revision"],
                                            "author": current_rev["author"],
                                        },
                                        citation=self._build_citation("codecortex-svn", position),
                                    ))
                                    position += 1
                                    if position > req.max_results:
                                        break
                                current_rev = {"revision": m.group(1), "author": m.group(2), "date": m.group(3).strip(), "message": ""}
                            elif current_rev is not None:
                                stripped = line.strip()
                                if stripped and not stripped.startswith("-----") and not stripped.startswith("Changed paths"):
                                    if current_rev["message"]:
                                        current_rev["message"] += " " + stripped
                                    else:
                                        current_rev["message"] = stripped
                except Exception as le:
                    logger.debug("search|svn|log_skipped=%s", str(le)[:80])

            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|svn|root=%s|results=%d|elapsed=%dms",
                         str(svn_root), len(results), elapsed)
            return results, {"results_found": len(results), "svn_root": str(svn_root)}, None
        except Exception as e:
            logger.warning("search|svn|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Provider 16: codecortex-blame (git blame & hotspots) ──────────
    async def _search_blame(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        try:
            from ..modules.filesystem.adapters.git import DiskGit
            from ..modules.codeanalysis.analyzers.audit import CodeAuditor

            root = req.repo_path or os.getcwd()
            root_path = Path(root).resolve()
            git_root = DiskGit.find_root(root_path)
            if not git_root:
                return [], {}, f"Not a git repository: {root}"

            results = []
            position = 1

            # Phase 1: Git history insights for repo
            if position <= req.max_results:
                try:
                    insights = DiskGit.get_insights(root_path)
                    if insights:
                        results.append(SearchResultItem(
                            title="[repo_insights] Repository Overview",
                            url="", display_url=str(git_root),
                            snippet=f"Branch: {insights.get('current_branch', '?')} | "
                                    f"Last author: {insights.get('last_author', '?')} | "
                                    f"Commits: {insights.get('commit_count', 0)} | "
                                    f"Lines: +{insights.get('lines_added', 0)}/-{insights.get('lines_deleted', 0)}",
                            position=position, score=0.95,
                            metadata={
                                "type": "repo_insights",
                                "source_type": "blame",
                                "file_path": str(git_root), "line": 0,
                                "branch": insights.get("current_branch", ""),
                                "last_author": insights.get("last_author", ""),
                                "last_commit": insights.get("last_commit_message", ""),
                            },
                            citation=self._build_citation("codecortex-blame", position),
                        ))
                        position += 1
                except Exception as ie:
                    logger.debug("search|blame|insights_skipped=%s", str(ie)[:80])

            # Phase 2: Git history scan for secrets/tags
            if position <= req.max_results:
                try:
                    history = CodeAuditor._scan_git_history(git_root)
                    for h in history:
                        if position > req.max_results:
                            break
                        results.append(SearchResultItem(
                            title=f"[git_history] {h.get('type', 'finding')}",
                            url="", display_url=f"git history line {h.get('line', 0)}",
                            snippet=f"{h.get('severity', 'medium')}: {h.get('type', 'finding')} found in git history",
                            position=position, score=0.8,
                            metadata={
                                "type": "git_history",
                                "source_type": "blame",
                                "file_path": "git history", "line": h.get("line", 0),
                                "priority": h.get("severity", "medium"),
                            },
                            citation=self._build_citation("codecortex-blame", position),
                        ))
                        position += 1
                except Exception as he:
                    logger.debug("search|blame|history_skipped=%s", str(he)[:80])

            # Phase 3: git blame on files matching query (or all tracked)
            if req.query and position <= req.max_results:
                try:
                    result = subprocess.run(
                        ["git", "-C", str(git_root), "ls-files"],
                        capture_output=True, text=True, timeout=30,
                    )
                    if result.returncode == 0:
                        query_lower = req.query.lower()
                        tracked = [l.strip() for l in result.stdout.splitlines() if l.strip()]
                        file_pats = [p.strip() for p in req.file_pattern.split(",")] if req.file_pattern else ["*"]
                        for tf in tracked:
                            if position > req.max_results:
                                break
                            if query_lower and query_lower not in tf.lower():
                                continue
                            if not any(fnmatch.fnmatch(tf, p) for p in file_pats):
                                continue
                            blame = subprocess.run(
                                ["git", "-C", str(git_root), "blame", "-L", "1,1", "--porcelain", "--", tf],
                                capture_output=True, text=True, timeout=15,
                            )
                            if blame.returncode != 0 or not blame.stdout.strip():
                                continue
                            blame_lines = blame.stdout.splitlines()
                            author = ""
                            commit = ""
                            for bl in blame_lines[:20]:
                                if bl.startswith("author "):
                                    author = bl[7:]
                                if bl.startswith("committer-time"):
                                    break
                            first_line = blame_lines[0] if blame_lines else ""
                            commit = first_line.split(" ")[0] if first_line else ""
                            results.append(SearchResultItem(
                                title=f"[blame] {tf}",
                                url=f"file://{git_root / tf}",
                                display_url=tf,
                                snippet=f"Last author: {author} | Commit: {commit[:12] if commit else '?'}",
                                position=position, score=0.75,
                                metadata={
                                    "type": "blame_info",
                                    "source_type": "blame",
                                    "file_path": tf, "line": 0,
                                    "last_author": author,
                                    "last_commit": commit[:12] if commit else "",
                                },
                                citation=self._build_citation("codecortex-blame", position),
                            ))
                            position += 1
                except Exception as be:
                    logger.debug("search|blame|file_scan_skipped=%s", str(be)[:80])

            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|blame|repo=%s|results=%d|elapsed=%dms",
                         str(git_root), len(results), elapsed)
            return results, {"results_found": len(results), "git_root": str(git_root)}, None
        except Exception as e:
            logger.warning("search|blame|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Provider 9: codecortex-agentart (strict <project>/.agents only) ─
    async def _search_agentart(self, req: SearchRequest) -> Tuple[List[SearchResultItem], Dict, Optional[str]]:
        t0 = time.monotonic()
        try:
            base = Path(req.repo_path) if req.repo_path else Path.cwd()
            _validate_path(str(base))
            agents_dir = base / ".agents"
            if not agents_dir.exists() or not agents_dir.is_dir():
                return [], {}, f"No .agents directory at: {agents_dir} (only <project>/.agents is supported)"

            results = []
            position = 1
            query_lower = req.query.lower() if req.query else ""
            allowed_patterns = ["*.md", "*.yml", "*.yaml", "*.json", "*.toml", "*.txt"]
            if req.artifact_type:
                allowed_patterns = [f"*.{req.artifact_type}"]

            for pat in allowed_patterns:
                for fpath in agents_dir.glob(f"**/{pat}"):
                    if position > req.max_results:
                        break
                    try:
                        fsize = fpath.stat().st_size
                        if fsize > 2 * 1024 * 1024:
                            continue
                        content = fpath.read_text(encoding="utf-8", errors="replace")
                    except Exception:
                        continue
                    rp = fpath.relative_to(base)
                    if query_lower and query_lower not in content[:10000].lower() and query_lower not in rp.name.lower():
                        continue
                    snippet = content[:500]
                    version_match = re.search(r'(?:version|rev|release)[:\s]*[\'"\s]*([\d][\w.\-+]*)',
                                              content[:500], re.IGNORECASE)
                    version_found = version_match.group(1) if version_match else None
                    st = fpath.stat()
                    results.append(SearchResultItem(
                        title=str(rp), url=f"file://{fpath}", display_url=str(rp),
                        snippet=snippet, position=position, score=0.65,
                        content=content[:2000],
                        published_at=datetime.fromtimestamp(st.st_mtime).isoformat(),
                        metadata={"type": "agent_artifact", "source_type": "agentart",
                                  "size_bytes": st.st_size, "artifact_version": version_found,
                                  "artifact_type": fpath.suffix.lstrip(".")},
                        citation=self._build_citation("codecortex-agentart", position),
                    ))
                    position += 1
                    if position > req.max_results:
                        break
                if position > req.max_results:
                    break
            elapsed = int((time.monotonic() - t0) * 1000)
            logger.info("search|agentart|query=%s|dir=%s|results=%d|elapsed=%dms",
                         req.query[:80], str(agents_dir), len(results), elapsed)
            return results, {"results_found": len(results), "scanned_dir": str(agents_dir)}, None
        except Exception as e:
            logger.warning("search|agentart|error=%s", str(e)[:100])
            return [], {}, str(e)

    # ── Orchestration (with auto-indexing) ─────────────────
    async def search(self, req: SearchRequest) -> SearchResponse:
        return await self._search_with_context(req, attempt=0)

    async def _search_with_context(self, req: SearchRequest, attempt: int = 0) -> SearchResponse:
        t0 = time.monotonic()
        provider = req.model
        errors: List[Dict[str, Any]] = []
        all_results: List[SearchResultItem] = []
        per_provider_metrics: Dict[str, Dict] = {}
        providers_used = 0
        index_status: Optional[Dict[str, Any]] = None

        provider_map = {
            "codecortex-combo": [
                "codebase", "repowt", "filesystem", "graph", "idegraph",
                "knowledge", "crossproject", "codeindex", "agentart", "codelogs",
                "todo", "stub", "security", "empty", "svn", "blame",
            ],
            "codecortex-codebase": ["codebase"],
            "codecortex-repowt": ["repowt"],
            "codecortex-filesystem": ["filesystem"],
            "codecortex-graph": ["graph"],
            "codecortex-idegraph": ["idegraph"],
            "codecortex-knowledge": ["knowledge"],
            "codecortex-crossproject": ["crossproject"],
            "codecortex-codeindex": ["codeindex"],
            "codecortex-agentart": ["agentart"],
            "codecortex-codelogs": ["codelogs"],
            "codecortex-todo": ["todo"],
            "codecortex-stub": ["stub"],
            "codecortex-svn": ["svn"],
            "codecortex-security": ["security"],
            "codecortex-empty": ["empty"],
            "codecortex-blame": ["blame"],
        }
        targets = provider_map.get(provider, ["codebase", "filesystem"])
        if req.search_type == "code":
            targets = ["codebase", "graph", "codeindex", "crossproject"]
        elif req.search_type == "file":
            targets = ["filesystem", "agentart"]
        elif req.search_type == "memory":
            targets = ["idegraph"]
        elif req.search_type == "knowledge":
            targets = ["knowledge"]
        elif req.search_type == "repo":
            targets = ["repowt"]
        elif req.search_type == "log":
            targets = ["codelogs"]
        elif req.search_type == "todo":
            targets = ["todo"]
        elif req.search_type == "stub":
            targets = ["stub"]
        elif req.search_type == "security":
            targets = ["security"]
        elif req.search_type == "empty":
            targets = ["empty"]
        elif req.search_type == "blame":
            targets = ["blame"]
        elif req.search_type == "svn":
            targets = ["svn"]

        search_fns: Dict[str, Any] = {
            "codebase": self._search_codebase,
            "repowt": self._search_repowt,
            "filesystem": self._search_filesystem,
            "graph": self._search_graph,
            "idegraph": self._search_idegraph,
            "knowledge": self._search_knowledge,
            "crossproject": self._search_crossproject,
            "codeindex": self._search_codeindex,
            "agentart": self._search_agentart,
            "codelogs": self._search_codelogs,
            "todo": self._search_todo,
            "stub": self._search_stub,
            "svn": self._search_svn,
            "security": self._search_security,
            "empty": self._search_empty,
            "blame": self._search_blame,
        }
        per_target_max = max(1, req.max_results // max(1, len(targets)))

        async def _run_one(target: str) -> Tuple[str, List[SearchResultItem], Dict, Optional[str]]:
            fn = search_fns.get(target)
            if not fn:
                return target, [], {}, f"Unknown target: {target}"
            sub_req = SearchRequest(
                query=req.query, model=target, max_results=per_target_max,
                search_type=req.search_type, repo_path=req.repo_path,
                repo_id=req.repo_id, symbol_type=req.symbol_type,
                language=req.language, file_pattern=req.file_pattern,
                content_regex=req.content_regex, recursive=req.recursive,
                max_depth=req.max_depth, search_mode=req.search_mode,
                project_name=req.project_name, ide_name=req.ide_name,
                knowledge_type=req.knowledge_type, direction=req.direction,
                relation_type=req.relation_type, graph_max_depth=req.graph_max_depth,
                status_filter=req.status_filter, commit_range=req.commit_range,
                diff_search=req.diff_search, since=req.since,
                min_references=req.min_references, include_signatures=req.include_signatures,
                artifact_type=req.artifact_type, version=req.version,
                include_history=req.include_history,
                auto_index=False,
            )
            t1 = time.monotonic()
            items, meta, err = await fn(sub_req)
            elapsed = int((time.monotonic() - t1) * 1000)
            if err:
                return target, [], {"latency_ms": elapsed, "error": err}, err
            return target, items, {"latency_ms": elapsed, "results": len(items), **meta}, None

        tasks = [_run_one(t) for t in targets]
        gathered = await asyncio.gather(*tasks, return_exceptions=True)

        for g in gathered:
            if isinstance(g, Exception):
                errors.append({"message": str(g), "provider": "unknown"})
                continue
            target, items, meta, err = g
            provider_key = f"codecortex-{target}"
            if err:
                errors.append({"message": err, "provider": provider_key})
            else:
                providers_used += 1
            per_provider_metrics[provider_key] = meta
            all_results.extend(items)

        # Check index status for auto-indexing
        db_providers = {"codebase", "graph", "crossproject", "codeindex"}
        active_dbs = [t for t in targets if t in db_providers]
        db_have_data = any(
            per_provider_metrics.get(f"codecortex-{t}", {}).get("results_found", 0) > 0
            for t in active_dbs
        ) if active_dbs else True

        needs_auto_index = (
            req.auto_index
            and attempt < 1
            and not db_have_data
            and req.repo_path is not None
            and os.path.isdir(req.repo_path)
        )

        if needs_auto_index:
            from ..services.auto_indexer import check_index_status, run_full_index
            status = check_index_status(self.db, req.repo_path)
            should_index = (
                not status.indexed
                or status.needs_update
                or req.force_update
                or req.regraph
                or req.reindex
            )
            if should_index:
                logger.info("search|auto_index|triggered|path=%s|indexed=%s|needs_update=%s",
                            req.repo_path, status.indexed, status.needs_update)
                try:
                    # Create orchestrator for indexing
                    from ..main import CortexOrchestrator
                    orch = CortexOrchestrator()
                    await run_full_index(orch, req.repo_path)
                    index_status = {"auto_indexed": True, "message": "Index rebuilt successfully"}
                except Exception as ie:
                    logger.warning("search|auto_index|failed=%s", str(ie)[:100])
                    index_status = {"auto_indexed": False, "error": str(ie)[:100]}

                # Re-run search after indexing
                return await self._search_with_context(req, attempt=attempt + 1)

        if index_status is None:
            from ..services.auto_indexer import check_index_status
            index_status = check_index_status(self.db, req.repo_path).to_dict()

        # Sort + filter + paginate
        all_results.sort(key=lambda r: r.score, reverse=True)
        all_results = self._apply_result_filter(all_results, req.result_filter)
        total_available = len(all_results)
        all_results = all_results[req.offset:req.offset + req.max_results]
        for i, r in enumerate(all_results, 1):
            r.position = i

        has_more = (req.offset + req.max_results) < total_available
        pagination = {
            "offset": req.offset,
            "limit": req.max_results,
            "total": total_available,
            "has_more": has_more,
            "next_offset": req.offset + req.max_results if has_more else None,
        }

        total_elapsed = int((time.monotonic() - t0) * 1000)
        logger.info("search|orchestrated|provider=%s|targets=%d|results=%d|errors=%d|elapsed=%dms",
                     provider, providers_used, len(all_results), len(errors), total_elapsed)

        return SearchResponse(
            provider=provider, query=req.query, results=all_results, answer=None,
            usage={"providers_used": providers_used, "total_results": len(all_results),
                   "total_available": total_available, "search_cost_usd": 0.0},
            metrics={"response_time_ms": total_elapsed, "per_provider": per_provider_metrics,
                     "index_status": index_status},
            errors=errors, pagination=pagination,
        )


# ────────────────────────────────────────────────────────────
_engine: Optional[UnifiedSearchEngine] = None

def get_search_engine(orchestrator: Any = None, db: Any = None) -> UnifiedSearchEngine:
    global _engine
    if _engine is None or orchestrator is not None:
        _engine = UnifiedSearchEngine(orchestrator=orchestrator, db=db)
    return _engine
